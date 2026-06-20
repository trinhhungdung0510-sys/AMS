#!/usr/bin/env python3
"""TIN NGHIA AMS Sales Deck — proportional 16:9 grid layout."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand (unchanged) ───────────────────────────────────────────────
GREEN = RGBColor(0x0B, 0x6B, 0x1B)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE = RGBColor(0xF3, 0x6A, 0x10)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
MID_GRAY = RGBColor(0xE5, 0xE7, 0xEB)
DARK = RGBColor(0x1F, 0x29, 0x37)

FONT = "Arial"
IMG_DIR = Path(__file__).parent / "images"
OUTPUT = Path(__file__).parent / "TIN_NGHIA_AMS_Sales_Deck.pptx"

# ── Proportional grid (16:9 = 13.333" × 7.5") ─────────────────────
# Margins ~5%, golden-ratio content split 38.2% / 61.8%
SW, SH = 13.333, 7.5
MX = 0.600          # horizontal margin
MY = 0.500          # vertical margin (cover / hero)
SB = 0.150          # sidebar width
CX = MX + SB        # content area left  = 0.750"
CW = SW - MX - CX   # content area width = 12.000"
CY = 1.380          # content area top (below header)
FOOT_Y = 7.080
CH = FOOT_Y - CY    # content area height = 5.700"

PHI = 0.6180339887
TEXT_W = round(CW * (1 - PHI), 3)          # 4.616"
GUT = 0.400
IMG_W = round(CW - TEXT_W - GUT, 3)        # 6.984"
IMG_H = round(IMG_W * 9 / 16, 3)           # 16:9 image height
IMG_Y = round(CY + (CH - IMG_H) / 2, 3)    # vertically centred

TEXT_X = CX
IMG_X = round(CX + TEXT_W + GUT, 3)
TEXT_X_REV = round(CX + IMG_W + GUT, 3)    # text-right when image is left
IMG_X_REV = CX
FRAME_PAD = 0.050
PANEL_W = round(SW * (1 - PHI), 3)           # cover / contact panel ≈ 5.09"

HDR_EYEBROW_Y = 0.280
HDR_TITLE_Y = 0.580
HDR_ACCENT_Y = 1.240
HDR_H = CY

GUT_COL = 0.333     # gutter between equal columns
COL4_W = round((CW - 3 * GUT_COL) / 4, 3)  # 4-column grid
COL3_W = round((CW - 2 * GUT_COL) / 3, 3)  # 3-column grid
COL2_W = round((CW - GUT_COL) / 2, 3)      # 2-column grid

SLIDE_W = Inches(SW)
SLIDE_H = Inches(SH)
_slide_num = 0


def I(v):
    return Inches(v)


# ── Primitives ──────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line=None, r=False):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if r else MSO_SHAPE.RECTANGLE, l, t, w, h
    )
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, spacing=1.0):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.name = FONT
    p.font.color.rgb = color
    p.alignment = align
    p.line_spacing = spacing
    return box


def img(slide, name, l, t, w, h):
    p = IMG_DIR / name
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)
        return True
    return False


def img_frame(slide, name, l, t, w, h, accent=GREEN):
    pad = FRAME_PAD
    rect(slide, I(l - pad), I(t - pad), I(w + pad * 2), I(h + pad * 2), WHITE, MID_GRAY, r=True)
    rect(slide, I(l - pad), I(t - pad), I(w + pad * 2), I(0.060), accent)
    img(slide, name, I(l), I(t), I(w), I(h))


def fit_image_in_zone(zone_w, zone_h):
    """16:9 image sized to fit inside zone without overflow."""
    h = min(zone_h - FRAME_PAD * 2, zone_w * 9 / 16)
    w = h * 16 / 9
    return round(w, 3), round(h, 3)


def bullets(slide, l, t, w, items, sz=16, color=DARK, gap=0.520):
    y = t
    for item in items:
        rect(slide, I(l), I(y + 0.080), I(0.120), I(0.120), ORANGE, r=True)
        txt(slide, I(l + 0.280), I(y), I(w - 0.280), I(0.450), item, sz=sz, color=color)
        y += gap


def eyebrow(slide, l, t, label, color=ORANGE):
    rect(slide, I(l), I(t), I(1.600), I(0.320), color, r=True)
    txt(slide, I(l + 0.120), I(t + 0.020), I(1.400), I(0.280),
        label.upper(), sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_shell(slide, title, section=None, num=None, sidebar=GREEN):
    rect(slide, I(0), I(0), I(SB), SLIDE_H, sidebar)
    rect(slide, I(SB), I(0), I(SW - SB), I(0.060), sidebar)

    if section:
        eyebrow(slide, CX, HDR_EYEBROW_Y, section)

    txt(slide, I(CX), I(HDR_TITLE_Y if section else 0.350),
        I(CW - 0.500), I(0.750), title, sz=26, bold=True, color=GREEN)
    rect(slide, I(CX), I(HDR_ACCENT_Y if section else 1.050), I(1.400), I(0.050), ORANGE)

    rect(slide, I(SB), I(FOOT_Y), I(SW - SB), I(SH - FOOT_Y), LIGHT_GRAY)
    txt(slide, I(CX), I(FOOT_Y + 0.060), I(1.500), I(0.300),
        "TIN NGHIA", sz=10, bold=True, color=GREEN)
    txt(slide, I(CX + 1.050), I(FOOT_Y + 0.060), I(0.800), I(0.300),
        "AMS", sz=10, bold=True, color=ORANGE)
    txt(slide, I(CX + 1.550), I(FOOT_Y + 0.060), I(8.000), I(0.300),
        "  ·  AI Giám sát An toàn Sinh học 24/7", sz=10, color=GRAY)
    if num:
        txt(slide, I(SW - MX - 0.500), I(FOOT_Y + 0.060), I(0.500), I(0.300),
            str(num).zfill(2), sz=10, bold=True, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)


def place_image(slide, name, accent=GREEN, x=IMG_X, y=IMG_Y, w=IMG_W, h=IMG_H):
    iw, ih = fit_image_in_zone(w, CH)
    iy = round(CY + (CH - ih) / 2, 3)
    ix = round(x + (w - iw) / 2, 3)
    img_frame(slide, name, ix, iy, iw, ih, accent)


def split_layout(slide, title, section, num, image, content_fn,
                 reverse=False, sidebar=GREEN, accent=GREEN):
    slide_shell(slide, title, section, num, sidebar)
    if reverse:
        ix, tx, ac = IMG_X_REV, TEXT_X_REV, ORANGE
    else:
        ix, tx, ac = IMG_X, TEXT_X, accent
    # Image first (background), text second — never overlap columns
    place_image(slide, image, ac, ix, IMG_Y, IMG_W, IMG_H)
    content_fn(slide, tx, CY, TEXT_W)


def stat_card(slide, l, t, w, h, number, label, accent=ORANGE):
    rect(slide, I(l), I(t), I(w), I(h), WHITE, accent, r=True)
    rect(slide, I(l), I(t), I(w), I(0.060), accent)
    txt(slide, I(l + 0.200), I(t + 0.180), I(w - 0.400), I(0.550),
        number, sz=32, bold=True, color=accent, align=PP_ALIGN.CENTER)
    txt(slide, I(l + 0.150), I(t + 0.820), I(w - 0.300), I(0.500),
        label, sz=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)


def numbered_row(slide, l, t, w, num, text, accent=ORANGE):
    rect(slide, I(l), I(t), I(0.420), I(0.420), accent, r=True)
    txt(slide, I(l), I(t + 0.040), I(0.420), I(0.380),
        str(num), sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, I(l + 0.550), I(t), I(w - 0.550), I(0.520), WHITE, MID_GRAY, r=True)
    txt(slide, I(l + 0.750), I(t + 0.100), I(w - 0.900), I(0.400),
        text, sz=15, bold=True, color=DARK)


def feature_row(slide, l, t, w, icon, text, accent=GREEN):
    rect(slide, I(l), I(t), I(0.500), I(0.500), accent, r=True)
    txt(slide, I(l), I(t + 0.060), I(0.500), I(0.400),
        icon, sz=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, I(l + 0.650), I(t + 0.080), I(w - 0.650), I(0.450),
        text, sz=17, bold=True, color=DARK)


def grid_x(cols, index):
    """Left edge of column `index` in a `cols`-column grid within content width."""
    col_w = (CW - (cols - 1) * GUT_COL) / cols
    return CX + index * (col_w + GUT_COL)


def grid_w(cols):
    return (CW - (cols - 1) * GUT_COL) / cols


# ── Slides ──────────────────────────────────────────────────────────
def s01(prs):
    global _slide_num
    _slide_num = 1
    s = blank(prs)
    bg(s)
    img(s, "slide01_cover.png", I(0), I(0), SLIDE_W, SLIDE_H)

    rect(s, I(0), I(0), I(PANEL_W + MX), SLIDE_H, WHITE)
    rect(s, I(PANEL_W + MX - 0.080), I(0), I(0.080), SLIDE_H, MID_GRAY)
    rect(s, I(0), I(0), I(SB), SLIDE_H, GREEN)

    card_x, card_y = MX, MY
    card_w = PANEL_W - MX
    card_h = SH - 2 * MY
    rect(s, I(card_x), I(card_y), I(card_w), I(card_h), WHITE, GREEN, r=True)
    rect(s, I(card_x), I(card_y), I(card_w), I(0.080), GREEN)

    px = card_x + 0.400
    txt(s, I(px), I(card_y + 0.450), I(card_w - 0.800), I(0.450),
        "TIN NGHIA", sz=20, bold=True, color=GREEN)
    txt(s, I(px), I(card_y + 0.900), I(card_w - 0.800), I(1.100),
        "AMS", sz=80, bold=True, color=ORANGE)
    txt(s, I(px), I(card_y + 2.100), I(card_w - 0.800), I(0.700),
        "AI GIÁM SÁT AN TOÀN\nSINH HỌC 24/7", sz=19, bold=True, color=DARK)
    rect(s, I(px), I(card_y + 2.950), I(2.200), I(0.050), ORANGE)
    txt(s, I(px), I(card_y + 3.200), I(card_w - 0.800), I(0.450),
        "Bảo vệ đàn heo – Bảo vệ lợi nhuận", sz=16, color=GRAY)
    rect(s, I(px), I(card_y + card_h - 1.200), I(3.000), I(0.620), ORANGE, r=True)
    txt(s, I(px), I(card_y + card_h - 1.080), I(3.000), I(0.450),
        "Khám phá ngay  →", sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    badge_x = PANEL_W + MX + GUT
    badge_w = 2.200
    badge_h = 0.550
    badge_gap = (CH - 3 * badge_h) / 2 + CY - MY
    for i, (label, color) in enumerate([("Camera AI", GREEN), ("Cloud AI", ORANGE), ("24/7", GREEN)]):
        y = badge_gap + i * (badge_h + badge_gap * 0.35)
        rect(s, I(badge_x), I(y), I(badge_w), I(badge_h), WHITE, color, r=True)
        txt(s, I(badge_x), I(y + 0.100), I(badge_w), I(0.400),
            label, sz=13, bold=True, color=color, align=PP_ALIGN.CENTER)


def s02(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s, LIGHT_GRAY)
    slide_shell(s, "THỰC TRẠNG AN TOÀN SINH HỌC TẠI TRẠI HEO", "Vấn đề", _slide_num, ORANGE)

    items = [
        "Giám sát thủ công",
        "Không thể theo dõi liên tục",
        "Dễ bỏ sót vi phạm",
        "Nguy cơ dịch bệnh",
    ]
    row_h = 0.900
    row_gap = (CH - len(items) * row_h) / (len(items) + 1)
    for i, item in enumerate(items):
        y = CY + row_gap + i * (row_h + row_gap)
        numbered_row(s, TEXT_X, y, TEXT_W, i + 1, item, ORANGE if i % 2 else GREEN)

    place_image(s, "slide02_manual_monitoring.png", ORANGE)


def s03(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    slide_shell(s, "THIỆT HẠI DO DỊCH BỆNH", "Vấn đề", _slide_num, ORANGE)

    hero_h = CH * 0.28
    rect(s, I(TEXT_X), I(CY), I(TEXT_W), I(hero_h), ORANGE, r=True)
    txt(s, I(TEXT_X + 0.200), I(CY + 0.150), I(TEXT_W - 0.400), I(hero_h - 0.300),
        "THIỆT HẠI CÓ THỂ\nLÊN ĐẾN HÀNG TỶ ĐỒNG",
        sz=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    damages = ["Giảm năng suất", "Tăng tỷ lệ chết", "Tăng chi phí điều trị", "Mất lợi nhuận"]
    card_w = (TEXT_W - GUT_COL) / 2
    card_h = (CH - hero_h - GUT_COL * 2) / 2
    base_y = CY + hero_h + GUT_COL
    for i, d in enumerate(damages):
        col, row = i % 2, i // 2
        x = TEXT_X + col * (card_w + GUT_COL)
        y = base_y + row * (card_h + GUT_COL)
        stat_card(s, x, y, card_w, card_h, "✕", d, ORANGE if col else GREEN)

    place_image(s, "slide03_disease_impact.png", ORANGE)


def s04(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    slide_shell(s, "CÁC HÀNH VI VI PHẠM AMS PHÁT HIỆN", "Giải pháp", _slide_num)

    row_h = CH * 0.22
    groups = [("01", "CON NGƯỜI", GREEN), ("02", "DI CHUYỂN", ORANGE),
              ("03", "PHƯƠNG TIỆN", GREEN), ("04", "ĐỘNG VẬT", ORANGE)]
    for i, (num, label, color) in enumerate(groups):
        x = grid_x(4, i)
        w = grid_w(4)
        rect(s, I(x), I(CY), I(w), I(row_h), WHITE, color, r=True)
        rect(s, I(x), I(CY), I(w), I(0.060), color)
        txt(s, I(x + 0.150), I(CY + 0.120), I(0.500), I(0.350),
            num, sz=11, bold=True, color=color)
        txt(s, I(x + 0.150), I(CY + row_h * 0.48), I(w - 0.300), I(0.450),
            label, sz=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    img_y = CY + row_h + GUT_COL
    img_h = CH - row_h - GUT_COL
    img_frame(s, "slide04_detection_groups.png", CX, img_y, CW, img_h, GREEN)


def _content_05(sl, l, t, w):
    txt(sl, I(l), I(t), I(w), I(0.350), "AMS phát hiện:", sz=13, bold=True, color=GRAY)
    bullets(sl, l, t + 0.450, w, [
        "Không tắm trước khi vào trại",
        "Không sát trùng tay / chân",
        "Sai màu quần áo",
        "Người lạ vào khu sản xuất",
        "Công nhân tiếp xúc người lạ",
    ], sz=15, gap=0.580)


def s05(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "CON NGƯỜI – HÀNH VI", "Phát hiện", _slide_num,
                 "slide05_human_behavior.png", _content_05)


def _content_06(sl, l, t, w):
    txt(sl, I(l), I(t), I(w), I(0.350), "AMS phát hiện:", sz=13, bold=True, color=GRAY)
    bullets(sl, l, t + 0.450, w, [
        "Đi từ vùng bẩn sang vùng sạch",
        "Đi sai tuyến ATSH",
        "Xâm nhập vùng cấm",
    ], sz=16, gap=0.720)
    note_h = 1.400
    note_y = t + CH - note_h - 0.200
    rect(sl, I(l), I(note_y), I(w), I(note_h), GREEN_LIGHT, GREEN, r=True)
    txt(sl, I(l + 0.200), I(note_y + 0.200), I(w - 0.400), I(note_h - 0.400),
        "Sơ đồ luồng di chuyển\nATSH trực quan trên camera",
        sz=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def s06(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "DI CHUYỂN – LUỒNG ĐI", "Phát hiện", _slide_num,
                 "slide06_movement_flow.png", _content_06, reverse=True, sidebar=ORANGE)


def _content_07(sl, l, t, w):
    txt(sl, I(l), I(t), I(w), I(0.350), "AMS phát hiện:", sz=13, bold=True, color=GRAY)
    bullets(sl, l, t + 0.450, w, [
        "Xe chưa sát trùng",
        "Sát trùng không đủ thời gian",
        "Công nhân tiếp xúc xe bắt heo",
        "Công nhân tiếp xúc xe cám",
    ], sz=15, gap=0.620)


def s07(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "XE CỘ – PHƯƠNG TIỆN", "Phát hiện", _slide_num,
                 "slide07_vehicles.png", _content_07)


def _content_08(sl, l, t, w):
    txt(sl, I(l), I(t), I(w), I(0.350), "AMS phát hiện:", sz=13, bold=True, color=GRAY)
    chip_w = (w - GUT_COL) / 2
    chip_h = 0.650
    chip_y0 = t + 0.500
    for i, a in enumerate(["Chó", "Mèo", "Chim", "Chuột"]):
        col, row = i % 2, i // 2
        x = l + col * (chip_w + GUT_COL)
        y = chip_y0 + row * (chip_h + GUT_COL * 0.8)
        rect(sl, I(x), I(y), I(chip_w), I(chip_h), WHITE, ORANGE if col else GREEN, r=True)
        txt(sl, I(x), I(y + 0.140), I(chip_w), I(0.400),
            a, sz=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    banner_h = 0.750
    banner_y = t + CH - banner_h - 0.150
    rect(sl, I(l), I(banner_y), I(w), I(banner_h), ORANGE, r=True)
    txt(sl, I(l + 0.150), I(banner_y + 0.150), I(w - 0.300), I(banner_h - 0.300),
        "Động vật là nguồn mang mầm bệnh nguy hiểm",
        sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def s08(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "ĐỘNG VẬT XÂM NHẬP", "Phát hiện", _slide_num,
                 "slide08_animal_intrusion.png", _content_08, reverse=True, sidebar=ORANGE)


def s09(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s, LIGHT_GRAY)
    slide_shell(s, "LỘ TRÌNH PHÁT TRIỂN AMS", "Tầm nhìn", _slide_num)

    col_w = grid_w(3)
    card_h = CH - 0.350
    card_y = CY + 0.350

    stages = [
        ("1", "GIAI ĐOẠN 1", "AMS AN TOÀN\nSINH HỌC",
         ["Giám sát ATSH", "Phát hiện vi phạm", "Cảnh báo thời gian thực"], GREEN, True, None),
        ("2", "GIAI ĐOẠN 2", "AMS THEO DÕI\nSỨC KHỎE HEO",
         ["Heo sốt · Heo ho", "Heo bất thường", "Heo chết"], ORANGE, False, None),
        ("3", "GIAI ĐOẠN 3", "AMS 360",
         ["Chấm điểm công nhân", "Đánh giá rủi ro", "Khuyến nghị vận hành", "Trợ lý AI"],
         GREEN, False, "AI THAY BẠN LÀM CHỦ TRẠI"),
    ]

    for i, (num, phase, name, items, color, active, sub) in enumerate(stages):
        x = grid_x(3, i)
        rect(s, I(x), I(card_y), I(col_w), I(card_h), WHITE, color, r=True)
        if active:
            rect(s, I(x + col_w * 0.28), I(card_y - 0.180), I(col_w * 0.44), I(0.320), ORANGE, r=True)
            txt(s, I(x + col_w * 0.28), I(card_y - 0.160), I(col_w * 0.44), I(0.280),
                "HIỆN TẠI", sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        hdr_h = 0.550
        rect(s, I(x), I(card_y), I(col_w), I(hdr_h), color)
        txt(s, I(x + 0.150), I(card_y + 0.100), I(col_w - 0.300), I(0.380),
            phase, sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, I(x + 0.150), I(card_y + hdr_h + 0.150), I(col_w - 0.300), I(0.850),
            name, sz=16, bold=True, color=color, align=PP_ALIGN.CENTER)
        if sub:
            txt(s, I(x + 0.150), I(card_y + hdr_h + 0.950), I(col_w - 0.300), I(0.350),
                sub, sz=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        list_y = card_y + hdr_h + (1.150 if sub else 0.900)
        bullets(s, x + 0.200, list_y, col_w - 0.400, items, sz=13, gap=0.460)


def _feat_content(sl, l, t, w, items):
    row_h = 0.900
    gap = (CH - len(items) * row_h) / (len(items) + 1)
    for i, item in enumerate(items):
        y = t + gap + i * (row_h + gap)
        feature_row(sl, l, y, w, str(i + 1), item, ORANGE if i % 2 else GREEN)


def s10(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "VẼ VÙNG ATSH TRÊN CAMERA", "Tính năng", _slide_num,
                 "slide10_zone_drawing.png",
                 lambda sl, l, t, w: _feat_content(sl, l, t, w, [
                     "Vùng sạch · Vùng bẩn · Vùng cấm",
                     "Vùng sát trùng · Vùng cách ly",
                     "Khách hàng tự vẽ — AMS tự giám sát",
                 ]))


def s11(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    split_layout(s, "PHÁT HIỆN VI PHẠM THỜI GIAN THỰC", "Tính năng", _slide_num,
                 "slide11_realtime_detection.png",
                 lambda sl, l, t, w: _feat_content(sl, l, t, w, [
                     "Phát hiện ngay lập tức",
                     "Chụp ảnh bằng chứng",
                     "Cảnh báo tức thời",
                 ]), reverse=True, sidebar=ORANGE)


def s12(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s, LIGHT_GRAY)
    slide_shell(s, "ẢNH VI PHẠM – BẰNG CHỨNG RÕ RÀNG", "Tính năng", _slide_num)

    card_h = CH * 0.28
    feats = [("📷", "Lưu ảnh vi phạm"), ("🔍", "Dễ truy xuất"), ("✓", "Dễ đánh giá")]
    card_w = (TEXT_W - 2 * GUT_COL) / 3
    for i, (icon, label) in enumerate(feats):
        x = TEXT_X + i * (card_w + GUT_COL)
        stat_card(s, x, CY, card_w, card_h, icon, label, GREEN if i != 1 else ORANGE)

    banner_h = CH * 0.18
    banner_y = CY + card_h + GUT_COL
    rect(s, I(TEXT_X), I(banner_y), I(TEXT_W), I(banner_h), ORANGE, r=True)
    txt(s, I(TEXT_X + 0.200), I(banner_y + 0.120), I(TEXT_W - 0.400), I(banner_h - 0.240),
        "Không video  ·  Chỉ lưu ảnh vi phạm",
        sz=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    place_image(s, "slide12_violation_evidence.png", GREEN)


def s13(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    slide_shell(s, "DASHBOARD QUẢN LÝ THÔNG MINH", "Tính năng", _slide_num)

    place_image(s, "slide13_dashboard.png", GREEN)

    metrics = [("Tổng vi phạm", GREEN), ("Theo camera", ORANGE),
               ("Theo khu vực", GREEN), ("Theo thời gian", ORANGE)]
    row_h = (CH - 3 * GUT_COL) / 4
    for i, (label, color) in enumerate(metrics):
        y = CY + i * (row_h + GUT_COL)
        rect(s, I(TEXT_X), I(y), I(TEXT_W), I(row_h), WHITE, color, r=True)
        rect(s, I(TEXT_X), I(y), I(0.080), I(row_h), color)
        txt(s, I(TEXT_X + 0.250), I(y + row_h * 0.28), I(TEXT_W - 0.400), I(row_h * 0.5),
            label, sz=16, bold=True, color=DARK)


def s14(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    slide_shell(s, "CẢNH BÁO THÔNG MINH", "Tính năng", _slide_num, ORANGE)

    txt(s, I(TEXT_X), I(CY), I(TEXT_W), I(0.350), "Gửi qua:", sz=13, bold=True, color=GRAY)

    chip_w = (TEXT_W - GUT_COL) / 2
    chip_h = (CH * 0.55 - GUT_COL) / 2
    base_y = CY + 0.450
    channels = [("Zalo", GREEN), ("Email", ORANGE), ("SMS", GREEN), ("App", ORANGE)]
    for i, (label, color) in enumerate(channels):
        col, row = i % 2, i // 2
        x = TEXT_X + col * (chip_w + GUT_COL)
        y = base_y + row * (chip_h + GUT_COL)
        rect(s, I(x), I(y), I(chip_w), I(chip_h), WHITE, color, r=True)
        rect(s, I(x), I(y), I(chip_w), I(0.060), color)
        txt(s, I(x), I(y + chip_h * 0.32), I(chip_w), I(0.500),
            label, sz=20, bold=True, color=color, align=PP_ALIGN.CENTER)

    hub_h = CH * 0.18
    hub_y = CY + CH - hub_h
    rect(s, I(TEXT_X), I(hub_y), I(TEXT_W), I(hub_h), ORANGE, r=True)
    txt(s, I(TEXT_X + 0.200), I(hub_y + hub_h * 0.22), I(TEXT_W - 0.400), I(hub_h * 0.6),
        "⚡  Cảnh báo tức thì qua mọi kênh", sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    place_image(s, "slide14_smart_alerts.png", ORANGE)


def s15(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)

    def fn(sl, l, t, w):
        row_h = 0.900
        gap = (CH - 3 * row_h - 0.900) / 4
        for i, item in enumerate(["Mã hóa dữ liệu", "Phân quyền truy cập", "Sao lưu tự động"]):
            y = t + gap + i * (row_h + gap)
            feature_row(sl, l, y, w, "🔒", item, GREEN)
        banner_h = 0.800
        banner_y = t + CH - banner_h - gap
        rect(sl, I(l), I(banner_y), I(w), I(banner_h), GREEN, r=True)
        txt(sl, I(l + 0.150), I(banner_y + 0.150), I(w - 0.300), I(banner_h - 0.300),
            "Dữ liệu trang trại là tài sản của khách hàng",
            sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    split_layout(s, "BẢO MẬT TUYỆT ĐỐI", "Tính năng", _slide_num,
                 "slide15_security.png", fn, reverse=True)


def s16(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s, LIGHT_GRAY)
    slide_shell(s, "LỢI ÍCH VƯỢT TRỘI", "Giá trị", _slide_num)

    card_w = (TEXT_W - GUT_COL) / 2
    card_h = (CH - GUT_COL) / 2
    benefits = [("24/7", "Giám sát liên tục", GREEN), ("0", "Không bỏ sót vi phạm", ORANGE),
                ("↓", "Giảm nguy cơ dịch bệnh", GREEN), ("↑", "Tăng tính tuân thủ", ORANGE)]
    for i, (icon, label, color) in enumerate(benefits):
        col, row = i % 2, i // 2
        x = TEXT_X + col * (card_w + GUT_COL)
        y = CY + row * (card_h + GUT_COL)
        stat_card(s, x, y, card_w, card_h, icon, label, color)

    place_image(s, "slide16_benefits.png", GREEN)


def s17(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    slide_shell(s, "PHÙ HỢP MỌI MÔ HÌNH TRANG TRẠI", "Giá trị", _slide_num)

    row_h = CH * 0.20
    models = [("S", "Trại nhỏ", GREEN), ("M", "Trại vừa", ORANGE),
              ("L", "Trại lớn", GREEN), ("+", "Nhiều trại", ORANGE)]
    for i, (icon, label, color) in enumerate(models):
        x = grid_x(4, i)
        w = grid_w(4)
        rect(s, I(x), I(CY), I(w), I(row_h), WHITE, color, r=True)
        txt(s, I(x), I(CY + 0.080), I(w), I(0.500),
            icon, sz=26, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, I(x), I(CY + row_h * 0.58), I(w), I(0.400),
            label, sz=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)

    img_y = CY + row_h + GUT_COL
    img_h = CH - row_h - GUT_COL
    img_frame(s, "slide17_farm_models.png", CX, img_y, CW, img_h, GREEN)


def s18(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    img(s, "slide18_guarantee.png", I(0), I(0), SLIDE_W, SLIDE_H)
    veil = rect(s, I(0), I(0), SLIDE_W, SLIDE_H, WHITE)
    veil.fill.transparency = 0.25

    card_w = CW * PHI
    card_h = CH * 0.75
    card_x = (SW - card_w) / 2
    card_y = CY + (CH - card_h) / 2
    rect(s, I(card_x), I(card_y), I(card_w), I(card_h), WHITE, GREEN, r=True)
    rect(s, I(card_x), I(card_y), I(card_w), I(0.100), GREEN)

    txt(s, I(card_x + 0.400), I(card_y + 0.450), I(card_w - 0.800), I(1.400),
        "ĐẢM BẢO AN TOÀN CHO\nTRANG TRẠI TỚI", sz=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(s, I(card_x + 0.400), I(card_y + 1.650), I(card_w - 0.800), I(1.000),
        "99,99%", sz=72, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    rect(s, I(card_x + card_w / 2 - 0.850), I(card_y + card_h - 1.650),
         I(1.700), I(0.060), ORANGE)
    txt(s, I(card_x + 0.400), I(card_y + card_h - 1.350), I(card_w - 0.800), I(0.550),
        "Sản xuất bền vững · Gia tăng lợi nhuận", sz=18, color=GRAY, align=PP_ALIGN.CENTER)

    txt(s, I(SW - MX - 0.500), I(FOOT_Y + 0.060), I(0.500), I(0.300),
        "18", sz=10, bold=True, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)


def s19(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s, LIGHT_GRAY)
    slide_shell(s, "BẢNG GIÁ DỊCH VỤ AMS", "Đầu tư", _slide_num)

    card_w = COL2_W
    card_h = CH - 0.720
    card_y = CY
    cta_h = 0.520
    cta_y = FOOT_Y - cta_h - 0.080

    plans = [
        ("AMS STANDARD", "1,5 – 2", "TRIỆU / THÁNG", GREEN, False, [
            "Tối đa 20 camera", "Hành vi vi phạm ATSH", "Dashboard quản lý",
            "Cảnh báo thời gian thực", "Báo cáo ATSH",
        ]),
        ("AMS PROFESSIONAL", "3 – 6", "TRIỆU / THÁNG", ORANGE, True, [
            "Trên 20 camera", "Hành vi vi phạm ATSH", "Dashboard nâng cao",
            "Báo cáo chuyên sâu", "Hỗ trợ ưu tiên",
        ]),
    ]

    for i, (name, price, unit, color, featured, items) in enumerate(plans):
        x = grid_x(2, i)
        rect(s, I(x), I(card_y), I(card_w), I(card_h), WHITE, color, r=True)
        hdr_h = 0.750
        rect(s, I(x), I(card_y), I(card_w), I(hdr_h), color)
        txt(s, I(x), I(card_y + 0.150), I(card_w), I(0.450),
            name, sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if featured:
            rect(s, I(x + card_w - 1.500), I(card_y - 0.180), I(1.400), I(0.360), ORANGE, r=True)
            txt(s, I(x + card_w - 1.500), I(card_y - 0.160), I(1.400), I(0.320),
                "PHỔ BIẾN", sz=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txt(s, I(x + 0.300), I(card_y + hdr_h + 0.200), I(card_w - 0.600), I(0.700),
            price, sz=36, bold=True, color=color, align=PP_ALIGN.CENTER)
        txt(s, I(x + 0.300), I(card_y + hdr_h + 0.850), I(card_w - 0.600), I(0.350),
            unit, sz=13, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
        rect(s, I(x + 0.500), I(card_y + hdr_h + 1.250), I(card_w - 1.000), I(0.020), MID_GRAY)
        bullets(s, x + 0.550, card_y + hdr_h + 1.400, card_w - 1.000, items, sz=14, gap=0.480)

    cta_w = 4.100
    rect(s, I((SW - cta_w) / 2), I(cta_y), I(cta_w), I(cta_h), ORANGE, r=True)
    txt(s, I((SW - cta_w) / 2), I(cta_y + 0.080), I(cta_w), I(cta_h - 0.160),
        "Liên hệ tư vấn ngay  →", sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def s20(prs):
    global _slide_num; _slide_num += 1
    s = blank(prs); bg(s)
    img(s, "slide20_contact.png", I(0), I(0), SLIDE_W, SLIDE_H)

    panel_right = PANEL_W + MX
    rect(s, I(0), I(0), I(panel_right), SLIDE_H, WHITE)
    rect(s, I(panel_right - 0.080), I(0), I(0.080), SLIDE_H, MID_GRAY)
    rect(s, I(0), I(0), I(SB), SLIDE_H, GREEN)

    px = MX + 0.100
    pw = PANEL_W - 0.200
    txt(s, I(px), I(MY + 0.400), I(pw), I(0.450),
        "TIN NGHIA", sz=22, bold=True, color=GREEN)
    txt(s, I(px), I(MY + 0.900), I(pw), I(1.000),
        "AMS", sz=64, bold=True, color=ORANGE)
    rect(s, I(px), I(MY + 2.050), I(2.000), I(0.050), ORANGE)
    txt(s, I(px), I(MY + 2.350), I(pw), I(0.700),
        "AI GIÁM SÁT AN TOÀN SINH HỌC 24/7", sz=14, bold=True, color=DARK)
    txt(s, I(px), I(MY + 2.950), I(pw), I(0.400),
        "Bảo vệ đàn heo – Bảo vệ lợi nhuận", sz=13, color=GRAY)

    info_h = SH - 2 * MY - 3.600
    info_y = MY + 3.500
    rect(s, I(px), I(info_y), I(pw), I(info_h), LIGHT_GRAY, GREEN, r=True)
    txt(s, I(px + 0.250), I(info_y + 0.250), I(pw - 0.500), I(0.700),
        "31 Hoa Viên, KĐT Đặng Xá\nThuận An, TP Hà Nội", sz=13, color=DARK)
    txt(s, I(px + 0.250), I(info_y + info_h - 0.900), I(pw - 0.500), I(0.650),
        "0964106955", sz=30, bold=True, color=ORANGE)

    cta_y = info_y + info_h + 0.350
    rect(s, I(px), I(cta_y), I(3.500), I(0.550), ORANGE, r=True)
    txt(s, I(px), I(cta_y + 0.080), I(3.500), I(0.420),
        "Liên hệ ngay  →", sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    global _slide_num
    _slide_num = 0
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for fn in [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
               s11, s12, s13, s14, s15, s16, s17, s18, s19, s20]:
        fn(prs)

    prs.save(str(OUTPUT))
    print(f"✓ Saved: {OUTPUT}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Grid: content {CW}\" × {CH}\" | split {TEXT_W}\" / {IMG_W}\" | image 16:9")


if __name__ == "__main__":
    main()
